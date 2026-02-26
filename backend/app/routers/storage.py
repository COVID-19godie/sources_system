from io import BytesIO
import mimetypes
from pathlib import Path
from urllib.parse import quote

from docx import Document
from fastapi import APIRouter, Depends, File, Form, HTTPException, Query, UploadFile, status
from fastapi.responses import StreamingResponse
from minio.error import S3Error
from openpyxl import load_workbook
from pptx import Presentation
from sqlalchemy.orm import Session

from app import models, schemas
from app.core.file_access_tokens import (
    FileAccessTokenError,
    build_storage_access_urls,
    decode_storage_file_token,
)
from app.core.html_preview import repair_html_preview
from app.core.office_config import build_office_config
from app.core.office_converter import (
    ensure_legacy_pdf_preview,
    ensure_presentation_pdf_preview,
    is_legacy_office_suffix,
)
from app.core import storage as storage_service
from app.core import rag_sync, trash_service
from app.core.config import settings
from app.deps import get_current_admin, get_current_user, get_db_read, get_db_write


router = APIRouter(tags=["storage"])

PREVIEW_MAX_BYTES = 20 * 1024 * 1024
VIDEO_EXTENSIONS = {".mp4", ".webm", ".mov", ".m4v", ".avi"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".gif"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".ogg", ".aac"}
WORD_EXTENSIONS = {".doc", ".docx"}
EXCEL_EXTENSIONS = {".xls", ".xlsx", ".csv"}
PPT_EXTENSIONS = {".ppt", ".pptx"}


def _s3_to_http_error(error: S3Error) -> HTTPException:
    code = (error.code or "").lower()
    if code in {"nosuchkey", "nosuchobject", "notfound"}:
        return HTTPException(status_code=404, detail="Object not found")
    return HTTPException(status_code=502, detail=f"Storage error: {error.code}")


def _detect_file_format(key: str) -> str:
    suffix = Path(key).suffix.lower()
    if suffix == ".md":
        return "markdown"
    if suffix in {".html", ".htm"}:
        return "html"
    if suffix == ".pdf":
        return "pdf"
    if suffix in VIDEO_EXTENSIONS:
        return "video"
    if suffix in IMAGE_EXTENSIONS:
        return "image"
    if suffix in AUDIO_EXTENSIONS:
        return "audio"
    if suffix in WORD_EXTENSIONS:
        return "word"
    if suffix in EXCEL_EXTENSIONS:
        return "excel"
    if suffix in PPT_EXTENSIONS:
        return "ppt"
    return "other"


def _is_office_mode(mode: str) -> bool:
    return mode in {"word", "excel", "ppt"}


def _resolve_legacy_preview_key(object_key: str) -> str | None:
    if not is_legacy_office_suffix(Path(object_key).suffix.lower()):
        return None
    try:
        preview_key = ensure_legacy_pdf_preview(object_key, force=False)
    except Exception:  # noqa: BLE001
        return None
    return preview_key or None


def _resolve_presentation_preview_key(object_key: str) -> str | None:
    suffix = Path(object_key).suffix.lower()
    if suffix not in {".ppt", ".pptx"}:
        return None
    try:
        preview_key = ensure_presentation_pdf_preview(object_key, force=False)
    except Exception:  # noqa: BLE001
        return None
    return preview_key or None


def _build_access_urls_for_key(
    object_key: str,
    *,
    user_id: int | None = None,
) -> schemas.AccessUrlsOut:
    open_url, download_url = build_storage_access_urls(object_key=object_key, user_id=user_id)
    return schemas.AccessUrlsOut(open_url=open_url, download_url=download_url)


def _content_disposition_header(filename: str, disposition: str) -> str:
    safe_name = filename.replace('"', "")
    encoded = quote(filename, safe="")
    return f"{disposition}; filename=\"{safe_name}\"; filename*=UTF-8''{encoded}"


def _name_from_key(key: str) -> str:
    return key.rstrip("/").split("/")[-1]


def _parent_prefix(prefix_or_key: str) -> str | None:
    clean = prefix_or_key.rstrip("/")
    if not clean:
        return None
    if "/" not in clean:
        return ""
    return f"{clean.rsplit('/', 1)[0]}/"


def _bootstrap_for_chapter(
    db: Session,
    chapter: models.Chapter,
    *,
    fail_on_no_sections: bool,
) -> tuple[schemas.StorageBootstrapOut, bool]:
    sections = (
        db.query(models.ResourceSection)
        .filter(
            models.ResourceSection.stage == chapter.stage,
            models.ResourceSection.subject == chapter.subject,
            models.ResourceSection.is_enabled.is_(True),
        )
        .order_by(models.ResourceSection.sort_order.asc(), models.ResourceSection.id.asc())
        .all()
    )
    if not sections:
        if fail_on_no_sections:
            raise HTTPException(status_code=400, detail="No enabled sections found for this chapter scope")
        return (
            schemas.StorageBootstrapOut(
                chapter_id=chapter.id,
                chapter_code=chapter.chapter_code,
                created_count=0,
                skipped_count=0,
                items=[],
            ),
            True,
        )

    items: list[schemas.StorageBootstrapItemOut] = []
    created_count = 0
    skipped_count = 0
    for section in sections:
        folder_key = storage_service.build_resource_object_prefix(
            chapter.chapter_code,
            section.code,
            chapter.volume_code,
        )
        try:
            if storage_service.object_exists(folder_key):
                skipped_count += 1
                items.append(
                    schemas.StorageBootstrapItemOut(
                        section_id=section.id,
                        section_code=section.code,
                        section_name=section.name,
                        folder_key=folder_key,
                        status="skipped",
                    )
                )
                continue

            storage_service.upload_bytes(
                b"",
                folder_key,
                content_type="application/x-directory",
            )
            created_count += 1
            items.append(
                schemas.StorageBootstrapItemOut(
                    section_id=section.id,
                    section_code=section.code,
                    section_name=section.name,
                    folder_key=folder_key,
                    status="created",
                )
            )
        except S3Error as error:
            raise _s3_to_http_error(error) from error

    return (
        schemas.StorageBootstrapOut(
            chapter_id=chapter.id,
            chapter_code=chapter.chapter_code,
            created_count=created_count,
            skipped_count=skipped_count,
            items=items,
        ),
        False,
    )


def _render_word_preview(data: bytes, suffix: str) -> str:
    if suffix != ".docx":
        return "当前仅支持 .docx 在线预览，.doc 请下载后查看。"

    doc = Document(BytesIO(data))
    lines = [item.text.strip() for item in doc.paragraphs if item.text.strip()]
    if not lines:
        return "文档可读取，但未提取到正文内容。"
    return "\n\n".join(lines[:200])


def _render_excel_preview(data: bytes, suffix: str) -> str:
    if suffix == ".csv":
        return data.decode("utf-8", errors="ignore")[:20000]

    if suffix != ".xlsx":
        return "当前仅支持 .xlsx/.csv 在线预览，.xls 请下载后查看。"

    workbook = load_workbook(BytesIO(data), read_only=True, data_only=True)
    chunks: list[str] = []
    try:
        for sheet in workbook.worksheets[:3]:
            chunks.append(f"## 工作表：{sheet.title}")
            row_count = 0
            for row in sheet.iter_rows(min_row=1, max_col=12, values_only=True):
                if row_count >= 40:
                    break
                values = [str(item) if item is not None else "" for item in row]
                if any(values):
                    chunks.append(" | ".join(values))
                    row_count += 1
            if row_count == 0:
                chunks.append("(空)")
            chunks.append("")
    finally:
        workbook.close()
    return "\n".join(chunks).strip() or "表格可读取，但未提取到可展示内容。"


def _render_ppt_preview(data: bytes, suffix: str) -> str:
    if suffix != ".pptx":
        return "当前仅支持 .pptx 在线预览，.ppt 请下载后查看。"

    presentation = Presentation(BytesIO(data))
    lines: list[str] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                slide_texts.append(text)
        if slide_texts:
            lines.append(f"## 第 {idx} 页")
            lines.extend(slide_texts[:8])
            lines.append("")
    return "\n".join(lines).strip() or "PPT 可读取，但未提取到文本内容。"



@router.get("/list", response_model=schemas.StorageListOut)
def list_objects(
    prefix: str = Query(default=""),
    _: models.User = Depends(get_current_user),
):
    normalized = storage_service.normalize_prefix(prefix)
    try:
        rows = storage_service.list_objects(prefix=normalized, recursive=False)
    except S3Error as error:
        raise _s3_to_http_error(error) from error

    items: list[schemas.StorageListItemOut] = []
    for row in rows:
        key = row.object_name
        if key == normalized:
            continue
        is_dir = key.endswith("/")
        name = _name_from_key(key)
        content_type = None if is_dir else (mimetypes.guess_type(name)[0] or None)
        items.append(
            schemas.StorageListItemOut(
                key=key,
                name=name,
                is_dir=is_dir,
                size=None if is_dir else int(getattr(row, "size", 0) or 0),
                updated_at=getattr(row, "last_modified", None),
                content_type=content_type,
            )
        )

    items.sort(key=lambda item: (not item.is_dir, item.name.lower()))
    return schemas.StorageListOut(
        prefix=normalized,
        parent_prefix=_parent_prefix(normalized),
        items=items,
    )


@router.post("/folder", response_model=schemas.StorageFolderOut, status_code=status.HTTP_201_CREATED)
def create_folder(
    payload: schemas.StorageCreateFolderRequest,
    _: models.User = Depends(get_current_admin),
):
    try:
        key = storage_service.create_folder(payload.prefix, payload.name)
    except ValueError as error:
        raise HTTPException(status_code=400, detail=str(error)) from error
    except S3Error as error:
        raise _s3_to_http_error(error) from error

    return schemas.StorageFolderOut(key=key)


@router.post("/bootstrap/chapter/{chapter_id}", response_model=schemas.StorageBootstrapOut)
def bootstrap_chapter_folders(
    chapter_id: int,
    db: Session = Depends(get_db_read),
    _: models.User = Depends(get_current_admin),
):
    chapter = db.query(models.Chapter).filter(models.Chapter.id == chapter_id).first()
    if not chapter:
        raise HTTPException(status_code=404, detail="Chapter not found")

    result, _ = _bootstrap_for_chapter(
        db,
        chapter,
        fail_on_no_sections=True,
    )
    return result


@router.post("/bootstrap/all-chapters", response_model=schemas.StorageBootstrapBatchOut)
def bootstrap_all_chapters_folders(
    stage: str | None = Query(default=None),
    subject: str | None = Query(default=None),
    enabled_only: bool = Query(default=True),
    db: Session = Depends(get_db_read),
    _: models.User = Depends(get_current_admin),
):
    query = db.query(models.Chapter)
    if stage:
        query = query.filter(models.Chapter.stage == stage)
    if subject:
        query = query.filter(models.Chapter.subject == subject)
    if enabled_only:
        query = query.filter(models.Chapter.is_enabled.is_(True))

    chapters = query.order_by(models.Chapter.grade.asc(), models.Chapter.chapter_code.asc()).all()
    if not chapters:
        raise HTTPException(status_code=404, detail="No chapters found")

    results: list[schemas.StorageBootstrapOut] = []
    no_section_chapter_ids: list[int] = []
    created_total = 0
    skipped_total = 0

    for chapter in chapters:
        result, no_sections = _bootstrap_for_chapter(
            db,
            chapter,
            fail_on_no_sections=False,
        )
        if no_sections:
            no_section_chapter_ids.append(chapter.id)
        created_total += result.created_count
        skipped_total += result.skipped_count
        results.append(result)

    return schemas.StorageBootstrapBatchOut(
        total_chapters=len(chapters),
        created_count=created_total,
        skipped_count=skipped_total,
        no_section_chapter_ids=no_section_chapter_ids,
        chapters=results,
    )


@router.post("/upload", response_model=schemas.StorageUploadOut, status_code=status.HTTP_201_CREATED)
def upload_file(
    prefix: str = Form(default=""),
    file: UploadFile = File(...),
    _: models.User = Depends(get_current_admin),
):
    try:
        key, file_size, content_type = storage_service.upload_file_to_prefix(file, prefix=prefix)
    except ValueError as error:
        raise HTTPException(status_code=400, detail=str(error)) from error
    except S3Error as error:
        raise _s3_to_http_error(error) from error

    try:
        if Path(key).suffix.lower() in {".ppt", ".pptx"}:
            ensure_presentation_pdf_preview(key, force=False)
        else:
            ensure_legacy_pdf_preview(key, force=False)
    except Exception:  # noqa: BLE001
        pass

    return schemas.StorageUploadOut(
        key=key,
        name=_name_from_key(key),
        size=file_size,
        content_type=content_type,
    )


@router.post("/rename", response_model=schemas.StorageRenameOut)
def rename_object(
    payload: schemas.StorageRenameRequest,
    _: models.User = Depends(get_current_admin),
):
    source_key = payload.source_key.strip().lstrip("/")
    if not source_key:
        raise HTTPException(status_code=400, detail="Invalid source key")

    target_name = payload.target_name.strip().strip("/")
    if not target_name or "/" in target_name:
        raise HTTPException(status_code=400, detail="target_name should not include '/'")

    is_dir = source_key.endswith("/")

    try:
        if is_dir:
            source_prefix = storage_service.normalize_prefix(source_key)
            parent = _parent_prefix(source_prefix) or ""
            target_prefix = f"{parent}{target_name}/"
            if source_prefix == target_prefix:
                return schemas.StorageRenameOut(key=target_prefix, moved_count=0)

            target_rows = storage_service.list_objects(prefix=target_prefix, recursive=True)
            if target_rows or storage_service.object_exists(target_prefix):
                raise HTTPException(status_code=409, detail="Target folder already exists")

            source_rows = storage_service.list_objects(prefix=source_prefix, recursive=True)
            has_placeholder = storage_service.object_exists(source_prefix)
            has_placeholder_in_rows = any(row.object_name == source_prefix for row in source_rows)
            if not source_rows and not has_placeholder:
                raise HTTPException(status_code=404, detail="Source folder not found")

            moved = 0
            for row in source_rows:
                old_key = row.object_name
                if not old_key.startswith(source_prefix):
                    continue
                suffix = old_key[len(source_prefix):]
                new_key = f"{target_prefix}{suffix}"
                storage_service.copy_object(old_key, new_key)
                moved += 1

            if has_placeholder and not has_placeholder_in_rows:
                storage_service.copy_object(source_prefix, target_prefix)
                moved += 1

            storage_service.delete_prefix(source_prefix)
            return schemas.StorageRenameOut(key=target_prefix, moved_count=moved)

        source = storage_service.normalize_key(source_key)
        if not storage_service.object_exists(source):
            raise HTTPException(status_code=404, detail="Source file not found")

        parent = _parent_prefix(source) or ""
        source_suffix = Path(source).suffix
        target_file_name = target_name
        if "." not in Path(target_name).name and source_suffix:
            target_file_name = f"{target_name}{source_suffix}"
        target = f"{parent}{target_file_name}"
        target = storage_service.normalize_key(target)

        if source == target:
            return schemas.StorageRenameOut(key=target, moved_count=0)
        if storage_service.object_exists(target):
            raise HTTPException(status_code=409, detail="Target file already exists")

        storage_service.copy_object(source, target)
        storage_service.delete_object(source)
        return schemas.StorageRenameOut(key=target, moved_count=1)
    except S3Error as error:
        raise _s3_to_http_error(error) from error


@router.delete("/object", response_model=schemas.StorageDeleteOut)
def delete_object(
    key: str = Query(..., min_length=1),
    db: Session = Depends(get_db_write),
    current_admin: models.User = Depends(get_current_admin),
):
    object_key = key.strip().lstrip("/")
    if not object_key:
        raise HTTPException(status_code=400, detail="Invalid key")

    try:
        if object_key.endswith("/"):
            normalized_prefix = storage_service.normalize_prefix(object_key)
            rows = storage_service.list_objects(prefix=normalized_prefix, recursive=True)
            has_placeholder = storage_service.object_exists(object_key)
            has_resource_records = (
                db.query(models.Resource)
                .filter(
                    models.Resource.object_key.like(f"{normalized_prefix}%"),
                    models.Resource.is_trashed.is_(False),
                )
                .first()
                is not None
            )
            if not rows and not has_placeholder and not has_resource_records:
                raise HTTPException(status_code=404, detail="Folder not found")

            if normalized_prefix.startswith("resources/"):
                trashed_items = trash_service.trash_storage_prefix(
                    db,
                    normalized_prefix,
                    source=trash_service.TRASH_SOURCE_STORAGE_API,
                    deleted_by=current_admin.id,
                )
                resource_ids = list(
                    {
                        item.resource_id
                        for item in trashed_items
                        if item.resource_id is not None
                    }
                )
                if resource_ids:
                    rag_sync.sync_resource_to_workspaces(
                        db,
                        resource_ids,
                        actor_id=current_admin.id,
                        reason="storage_delete_prefix",
                    )
                db.commit()
                return schemas.StorageDeleteOut(
                    deleted_count=0,
                    trashed_count=len(trashed_items),
                    trashed_resource_count=sum(1 for item in trashed_items if item.resource_id is not None),
                    trashed_storage_count=sum(1 for item in trashed_items if item.resource_id is None),
                )

            deleted = storage_service.delete_prefix(object_key)
            return schemas.StorageDeleteOut(deleted_count=deleted)

        normalized = storage_service.normalize_key(object_key)
        has_object = storage_service.object_exists(normalized)
        has_resource = (
            db.query(models.Resource)
            .filter(models.Resource.object_key == normalized, models.Resource.is_trashed.is_(False))
            .first()
            is not None
        )
        if not has_object and not has_resource:
            raise HTTPException(status_code=404, detail="Object not found")

        if normalized.startswith("resources/"):
            item = trash_service.trash_storage_object(
                db,
                normalized,
                source=trash_service.TRASH_SOURCE_STORAGE_API,
                deleted_by=current_admin.id,
                meta={"from": "storage.delete"},
            )
            if item is None:
                raise HTTPException(status_code=404, detail="Object not found")
            if item.resource_id is not None:
                rag_sync.sync_resource_to_workspaces(
                    db,
                    [item.resource_id],
                    actor_id=current_admin.id,
                    reason="storage_delete_object",
                )
            db.commit()
            return schemas.StorageDeleteOut(
                deleted_count=0,
                trashed_count=1,
                trashed_resource_count=1 if item.resource_id is not None else 0,
                trashed_storage_count=0 if item.resource_id is not None else 1,
            )

        storage_service.delete_object(normalized)
        return schemas.StorageDeleteOut(deleted_count=1)
    except S3Error as error:
        raise _s3_to_http_error(error) from error


@router.post("/reconcile", response_model=schemas.StorageReconcileOut)
def reconcile_storage(
    dry_run: bool = Query(default=False),
    db: Session = Depends(get_db_write),
    current_admin: models.User = Depends(get_current_admin),
):
    try:
        result = trash_service.reconcile_missing_resources(db, dry_run=dry_run)
        resource_ids = [int(item) for item in (result.get("resource_ids") or [])]
        if not dry_run and resource_ids:
            rag_sync.sync_resource_to_workspaces(
                db,
                resource_ids,
                actor_id=current_admin.id,
                reason="storage_reconcile",
            )
        if dry_run:
            db.rollback()
        else:
            db.commit()
        return schemas.StorageReconcileOut(
            scanned_count=result["scanned_count"],
            trashed_count=result["trashed_count"],
            missing_count=result["missing_count"],
            dry_run=dry_run,
        )
    except S3Error as error:
        raise _s3_to_http_error(error) from error


@router.get("/download-url", response_model=schemas.StorageDownloadOut)
def download_url(
    key: str = Query(..., min_length=1),
    current_user: models.User = Depends(get_current_user),
):
    object_key = key.strip().lstrip("/")
    if object_key.endswith("/"):
        raise HTTPException(status_code=400, detail="Folder has no download URL")

    try:
        normalized = storage_service.normalize_key(object_key)
        if not storage_service.object_exists(normalized):
            raise HTTPException(status_code=404, detail="Object not found")
        access = _build_access_urls_for_key(normalized, user_id=current_user.id)
    except ValueError as error:
        raise HTTPException(status_code=400, detail=str(error)) from error
    except S3Error as error:
        raise _s3_to_http_error(error) from error

    return schemas.StorageDownloadOut(key=normalized, url=access.download_url)


@router.get("/access-urls", response_model=schemas.AccessUrlsOut)
def access_urls(
    key: str = Query(..., min_length=1),
    current_user: models.User = Depends(get_current_user),
):
    object_key = key.strip().lstrip("/")
    if object_key.endswith("/"):
        raise HTTPException(status_code=400, detail="Folder has no access URL")

    try:
        normalized = storage_service.normalize_key(object_key)
        if not storage_service.object_exists(normalized):
            raise HTTPException(status_code=404, detail="Object not found")
        return _build_access_urls_for_key(normalized, user_id=current_user.id)
    except ValueError as error:
        raise HTTPException(status_code=400, detail=str(error)) from error
    except S3Error as error:
        raise _s3_to_http_error(error) from error


@router.api_route("/file/{token}", methods=["GET", "HEAD"])
def stream_file(token: str):
    try:
        payload = decode_storage_file_token(token)
        object_key = storage_service.normalize_key(str(payload["obj"]))
        disposition = str(payload.get("disp") or "inline")
    except FileAccessTokenError as error:
        raise HTTPException(status_code=401, detail=str(error)) from error
    except ValueError as error:
        raise HTTPException(status_code=400, detail=str(error)) from error

    try:
        stat = storage_service.stat_object(object_key)
        response = storage_service.get_object_stream(object_key)
    except S3Error as error:
        raise _s3_to_http_error(error) from error

    filename = _name_from_key(object_key)
    content_type = (
        getattr(stat, "content_type", None)
        or mimetypes.guess_type(filename)[0]
        or "application/octet-stream"
    )
    headers = {
        "Content-Disposition": _content_disposition_header(filename, disposition),
        "Cache-Control": "private, max-age=60",
    }

    def iterator():
        try:
            while True:
                chunk = response.read(1024 * 1024)
                if not chunk:
                    break
                yield chunk
        finally:
            response.close()
            response.release_conn()

    return StreamingResponse(iterator(), media_type=content_type, headers=headers)


@router.get("/office-config", response_model=schemas.OfficeConfigOut)
def office_config(
    key: str = Query(..., min_length=1),
    current_user: models.User = Depends(get_current_user),
):
    if not settings.ONLYOFFICE_ENABLED:
        raise HTTPException(status_code=503, detail="OnlyOffice is disabled")

    object_key = key.strip().lstrip("/")
    if not object_key or object_key.endswith("/"):
        raise HTTPException(status_code=400, detail="File key is required")

    normalized = storage_service.normalize_key(object_key)
    mode = _detect_file_format(normalized)
    if not _is_office_mode(mode):
        raise HTTPException(status_code=400, detail="Object is not an office file")
    if is_legacy_office_suffix(Path(normalized).suffix.lower()):
        raise HTTPException(status_code=400, detail="Legacy office format uses PDF preview")
    if not storage_service.object_exists(normalized):
        raise HTTPException(status_code=404, detail="Object not found")

    return build_office_config(
        object_key=normalized,
        filename=_name_from_key(normalized),
        current_user=current_user,
        editable=current_user.role == models.UserRole.admin,
    )


@router.get("/preview", response_model=schemas.StoragePreviewOut)
def preview_object(
    key: str = Query(..., min_length=1),
    current_user: models.User = Depends(get_current_user),
):
    object_key = key.strip().lstrip("/")
    if object_key.endswith("/"):
        raise HTTPException(status_code=400, detail="Folder cannot be previewed")

    try:
        normalized = storage_service.normalize_key(object_key)
        stat = storage_service.stat_object(normalized)
        mode = _detect_file_format(normalized)
        content_type = getattr(stat, "content_type", None)
        file_size = int(getattr(stat, "size", 0) or 0)
        access = _build_access_urls_for_key(normalized, user_id=current_user.id)

        if mode == "ppt":
            presentation_pdf_key = _resolve_presentation_preview_key(normalized)
            if presentation_pdf_key:
                presentation_access = _build_access_urls_for_key(presentation_pdf_key, user_id=current_user.id)
                return schemas.StoragePreviewOut(
                    key=normalized,
                    mode="pdf",
                    content_type="application/pdf",
                    size=file_size,
                    url=presentation_access.open_url,
                    open_url=presentation_access.open_url,
                    download_url=presentation_access.download_url,
                )
            if is_legacy_office_suffix(Path(normalized).suffix.lower()):
                legacy_pdf_key = _resolve_legacy_preview_key(normalized)
                if legacy_pdf_key:
                    legacy_access = _build_access_urls_for_key(legacy_pdf_key, user_id=current_user.id)
                    return schemas.StoragePreviewOut(
                        key=normalized,
                        mode="pdf",
                        content_type="application/pdf",
                        size=file_size,
                        url=legacy_access.open_url,
                        open_url=legacy_access.open_url,
                        download_url=legacy_access.download_url,
                    )
            return schemas.StoragePreviewOut(
                key=normalized,
                mode="other",
                content_type=content_type,
                size=file_size,
                url=access.open_url,
                open_url=access.open_url,
                download_url=access.download_url,
                content="PPT 预览转换失败，请稍后重试或下载查看。",
            )

        if _is_office_mode(mode):
            return schemas.StoragePreviewOut(
                key=normalized,
                mode=mode,
                content_type=content_type,
                size=file_size,
                url=access.open_url,
                open_url=access.open_url,
                download_url=access.download_url,
            )

        if mode in {"pdf", "video", "image", "audio"}:
            return schemas.StoragePreviewOut(
                key=normalized,
                mode=mode,
                content_type=content_type,
                size=file_size,
                url=access.open_url,
                open_url=access.open_url,
                download_url=access.download_url,
            )

        if mode == "other":
            return schemas.StoragePreviewOut(
                key=normalized,
                mode=mode,
                content_type=content_type,
                size=file_size,
                url=access.open_url,
                open_url=access.open_url,
                download_url=access.download_url,
            )

        if mode == "html":
            repaired_html: str | None = None
            try:
                data = storage_service.get_object_bytes(normalized, max_bytes=PREVIEW_MAX_BYTES)
                repaired_html = repair_html_preview(data.decode("utf-8", errors="ignore"))
            except ValueError:
                repaired_html = None
            except Exception:  # noqa: BLE001
                repaired_html = None
            return schemas.StoragePreviewOut(
                key=normalized,
                mode=mode,
                content_type=content_type,
                size=file_size,
                url=access.open_url,
                open_url=access.open_url,
                download_url=access.download_url,
                content=repaired_html,
            )

        data = storage_service.get_object_bytes(normalized, max_bytes=PREVIEW_MAX_BYTES)
        suffix = Path(normalized).suffix.lower()

        try:
            if mode == "markdown":
                content = data.decode("utf-8", errors="ignore")
            elif mode == "word":
                content = _render_word_preview(data, suffix)
            elif mode == "excel":
                content = _render_excel_preview(data, suffix)
            elif mode == "ppt":
                content = _render_ppt_preview(data, suffix)
            else:
                content = ""
        except Exception:  # noqa: BLE001
            content = "预览解析失败，请下载后查看完整内容。"

        return schemas.StoragePreviewOut(
            key=normalized,
            mode=mode,
            content_type=content_type,
            size=file_size,
            open_url=access.open_url,
            download_url=access.download_url,
            content=content,
        )
    except ValueError as error:
        raise HTTPException(status_code=400, detail=str(error)) from error
    except S3Error as error:
        raise _s3_to_http_error(error) from error
