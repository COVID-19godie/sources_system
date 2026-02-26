import logging
import shutil
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
import re
from urllib.parse import urlparse
from uuid import uuid4

from docx import Document
from fastapi import APIRouter, BackgroundTasks, Depends, File, Form, HTTPException, Query, UploadFile, status
from openpyxl import load_workbook
from pptx import Presentation
from sqlalchemy import func, or_, select, text
from sqlalchemy.orm import Session

from app import models, schemas
from app.core import ai_service, chapter_classifier, rag_sync, resource_variants
from app.core import semantic_ranker
from app.core.config import settings
from app.core.file_access_tokens import build_storage_access_urls
from app.core.html_preview import repair_html_preview
from app.core.office_config import build_office_config
from app.core.office_converter import (
    ensure_legacy_pdf_preview,
    ensure_presentation_pdf_preview,
    is_legacy_office_suffix,
)
from app.core import trash_service
from app.core.mineru_api import MinerUAPIError, text_to_markdown_with_mineru
from app.core.storage import (
    build_resource_object_key,
    get_object_bytes,
    upload_file,
)
from app.deps import (
    get_auth_payload_optional,
    get_current_admin,
    get_current_user,
    get_db_read,
    get_db_write,
)


router = APIRouter(tags=["resources"])
logger = logging.getLogger(__name__)
UPLOAD_DIR = Path("uploads")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

VIDEO_EXTENSIONS = {".mp4", ".webm", ".mov", ".m4v", ".avi"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".gif"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".ogg", ".aac"}
WORD_EXTENSIONS = {".doc", ".docx"}
EXCEL_EXTENSIONS = {".xls", ".xlsx", ".csv"}
PPT_EXTENSIONS = {".ppt", ".pptx"}
PREVIEW_MAX_BYTES = 20 * 1024 * 1024


def parse_tags(raw_tags: str | None) -> list[str]:
    if not raw_tags:
        return []
    cleaned = [tag.strip() for tag in raw_tags.split(",") if tag.strip()]
    deduped = list(dict.fromkeys(cleaned))
    return deduped


def parse_tag_list(raw_tags: list[str] | None) -> list[str]:
    if not raw_tags:
        return []
    cleaned = [str(tag).strip() for tag in raw_tags if str(tag).strip()]
    return list(dict.fromkeys(cleaned))


def compose_resource_name(
    raw_title: str | None,
    chapter: models.Chapter | None,
    section: models.ResourceSection | None,
    filename: str | None,
    tags: list[str] | None = None,
    description: str | None = None,
    volume_code_override: str | None = None,
) -> tuple[str, str, bool]:
    fallback_keyword = chapter_classifier.clean_filename_stem(filename, fallback="资源")
    rich_keyword = chapter_classifier.normalize_keyword(
        " ".join(
            item
            for item in [
                " ".join(tags or []),
                (raw_title or "").strip(),
                (description or "").strip(),
                Path(filename or "").stem,
            ]
            if item
        ),
        fallback=fallback_keyword,
    )
    chapter_code = chapter.chapter_code if chapter else None
    volume_code = chapter.volume_code if chapter else (volume_code_override or None)
    section_code = section.code if section else None

    if raw_title and raw_title.strip():
        custom_title = raw_title.strip()[:255]
        base_name = chapter_classifier.normalize_keyword(custom_title, fallback=fallback_keyword)
        return custom_title, base_name, False

    composed = chapter_classifier.build_resource_title(
        volume_code=volume_code,
        chapter_code=chapter_code,
        section_code=section_code,
        keyword=rich_keyword,
    )
    return composed, composed, True


def detect_file_format(filename: str | None) -> str:
    if not filename:
        return "other"

    suffix = Path(filename).suffix.lower()
    if suffix == ".md":
        return "markdown"
    if suffix in {".html", ".htm"}:
        return "html"
    if suffix == ".pdf":
        return "pdf"
    if suffix in VIDEO_EXTENSIONS:
        return "video"
    if suffix in IMAGE_EXTENSIONS:
        return "image"
    if suffix in AUDIO_EXTENSIONS:
        return "audio"
    if suffix in WORD_EXTENSIONS:
        return "word"
    if suffix in EXCEL_EXTENSIONS:
        return "excel"
    if suffix in PPT_EXTENSIONS:
        return "ppt"
    return "other"


def _read_upload_preview_text(file: UploadFile | None, file_format: str) -> str:
    if file is None:
        return ""
    if file_format not in {"markdown", "html", "pdf", "word", "excel", "ppt"}:
        return ""
    try:
        data = file.file.read(PREVIEW_MAX_BYTES + 1)
        file.file.seek(0)
        if len(data) > PREVIEW_MAX_BYTES:
            data = data[:PREVIEW_MAX_BYTES]
    except Exception:  # noqa: BLE001
        return ""

    suffix = Path(file.filename or "").suffix.lower()
    try:
        if file_format == "markdown":
            return data.decode("utf-8", errors="ignore")[:4000]
        if file_format == "html":
            return strip_html(data.decode("utf-8", errors="ignore"))[:4000]
        if file_format == "pdf":
            return render_pdf_text(data)[:4000]
        if file_format == "word":
            return render_word_preview(data, suffix)[:4000]
        if file_format == "excel":
            return render_excel_preview(data, suffix)[:4000]
        if file_format == "ppt":
            return render_ppt_preview(data, suffix)[:4000]
    except Exception:  # noqa: BLE001
        return ""
    return ""


def resolve_chapter_ids(resource: models.Resource) -> list[int]:
    ids = set()
    if resource.chapter_id:
        ids.add(resource.chapter_id)
    for link in resource.chapter_links:
        ids.add(link.chapter_id)
    return sorted(ids)


def resolve_access_urls(
    resource: models.Resource,
    *,
    user_id: int | None = None,
) -> schemas.AccessUrlsOut | None:
    if resource.storage_provider == models.StorageProvider.minio and resource.object_key:
        try:
            open_url, download_url = build_storage_access_urls(
                object_key=resource.object_key,
                user_id=user_id,
            )
            return schemas.AccessUrlsOut(open_url=open_url, download_url=download_url)
        except Exception:  # noqa: BLE001
            return None

    if resource.file_path:
        return schemas.AccessUrlsOut(open_url=resource.file_path, download_url=resource.file_path)

    return None


def resolve_download_url(resource: models.Resource) -> str | None:
    access = resolve_access_urls(resource)
    if not access:
        return None
    return access.download_url


def resolve_resource_filename(resource: models.Resource) -> str:
    if resource.object_key:
        return Path(resource.object_key).name
    if resource.file_path:
        return Path(resource.file_path).name
    return ""


def resolve_resource_suffix(resource: models.Resource) -> str:
    return Path(resolve_resource_filename(resource)).suffix.lower()


def is_legacy_office_resource(resource: models.Resource) -> bool:
    return is_legacy_office_suffix(resolve_resource_suffix(resource))


def is_office_resource(resource: models.Resource) -> bool:
    return (resource.file_format or "other") in {"word", "excel", "ppt"}


def resolve_legacy_preview_key(resource: models.Resource) -> str | None:
    if resource.storage_provider != models.StorageProvider.minio or not resource.object_key:
        return None
    if not is_legacy_office_resource(resource):
        return None
    try:
        preview_key = ensure_legacy_pdf_preview(resource.object_key, force=False)
    except Exception:  # noqa: BLE001
        return None
    return preview_key or None


def resolve_presentation_preview_key(resource: models.Resource) -> str | None:
    if resource.storage_provider != models.StorageProvider.minio or not resource.object_key:
        return None
    if (resource.file_format or "other") != "ppt":
        return None
    try:
        preview_key = ensure_presentation_pdf_preview(resource.object_key, force=False)
    except Exception:  # noqa: BLE001
        return None
    return preview_key or None


def read_preview_bytes(resource: models.Resource) -> bytes:
    if resource.storage_provider == models.StorageProvider.minio and resource.object_key:
        try:
            return get_object_bytes(resource.object_key, max_bytes=PREVIEW_MAX_BYTES)
        except ValueError as error:
            raise HTTPException(status_code=400, detail=str(error)) from error

    if resource.file_path:
        filename = resource.file_path.split("/")[-1]
        local_path = UPLOAD_DIR / filename
        if not local_path.exists():
            raise HTTPException(status_code=404, detail="Preview file not found")
        size = local_path.stat().st_size
        if size > PREVIEW_MAX_BYTES:
            raise HTTPException(status_code=400, detail="File too large for preview")
        return local_path.read_bytes()

    raise HTTPException(status_code=404, detail="Preview file not found")


def render_word_preview(data: bytes, suffix: str) -> str:
    if suffix != ".docx":
        return "当前仅支持 .docx 在线预览，.doc 请下载后查看。"

    doc = Document(BytesIO(data))
    lines = [item.text.strip() for item in doc.paragraphs if item.text.strip()]
    if not lines:
        return "文档可读取，但未提取到正文内容。"
    return "\n\n".join(lines[:200])


def render_excel_preview(data: bytes, suffix: str) -> str:
    if suffix == ".csv":
        return data.decode("utf-8", errors="ignore")[:20000]

    if suffix != ".xlsx":
        return "当前仅支持 .xlsx/.csv 在线预览，.xls 请下载后查看。"

    workbook = load_workbook(BytesIO(data), read_only=True, data_only=True)
    chunks: list[str] = []
    try:
        for sheet in workbook.worksheets[:3]:
            chunks.append(f"## 工作表：{sheet.title}")
            row_count = 0
            for row in sheet.iter_rows(min_row=1, max_col=12, values_only=True):
                if row_count >= 40:
                    break
                values = [str(item) if item is not None else "" for item in row]
                if any(values):
                    chunks.append(" | ".join(values))
                    row_count += 1
            if row_count == 0:
                chunks.append("(空)")
            chunks.append("")
    finally:
        workbook.close()

    return "\n".join(chunks).strip() or "表格可读取，但未提取到可展示内容。"


def render_ppt_preview(data: bytes, suffix: str) -> str:
    if suffix != ".pptx":
        return "当前仅支持 .pptx 在线预览，.ppt 请下载后查看。"

    presentation = Presentation(BytesIO(data))
    lines: list[str] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                slide_texts.append(text)
        if slide_texts:
            lines.append(f"## 第 {idx} 页")
            lines.extend(slide_texts[:8])
            lines.append("")

    return "\n".join(lines).strip() or "PPT 可读取，但未提取到文本内容。"



def render_pdf_text(data: bytes) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception:  # noqa: BLE001
        return ""

    try:
        reader = PdfReader(BytesIO(data))
    except Exception:  # noqa: BLE001
        return ""

    lines: list[str] = []
    for page in reader.pages[:20]:
        try:
            text = page.extract_text() or ""
        except Exception:  # noqa: BLE001
            text = ""
        if text.strip():
            lines.append(text.strip())
    return "\n\n".join(lines)


def strip_html(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r"<script[\s\S]*?</script>", " ", text, flags=re.IGNORECASE)
    text = re.sub(r"<style[\s\S]*?</style>", " ", text, flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def extract_text_for_ai(resource: models.Resource) -> str:
    parts: list[str] = []
    if resource.title:
        parts.append(resource.title)
    if resource.description:
        parts.append(resource.description)
    if resource.tags:
        parts.append("标签：" + "、".join(resource.tags))

    mode = resource.file_format or "other"
    try:
        if mode in {"markdown", "html", "word", "excel", "ppt", "pdf"}:
            raw = read_preview_bytes(resource)
            suffix = resolve_resource_suffix(resource)
            if mode == "markdown":
                parts.append(raw.decode("utf-8", errors="ignore"))
            elif mode == "html":
                parts.append(strip_html(raw.decode("utf-8", errors="ignore")))
            elif mode == "word":
                parts.append(render_word_preview(raw, suffix))
            elif mode == "excel":
                parts.append(render_excel_preview(raw, suffix))
            elif mode == "ppt":
                parts.append(render_ppt_preview(raw, suffix))
            else:
                parts.append(render_pdf_text(raw))
    except Exception:  # noqa: BLE001
        pass

    text = "\n\n".join(item for item in parts if item and item.strip())
    return text[: settings.AI_MAX_SOURCE_CHARS]


def _update_embedding_vector_column(db: Session, *, resource_id: int, embedding: list[float] | None) -> None:
    if not settings.SEMANTIC_PGVECTOR_ENABLED:
        return
    if not embedding:
        return
    try:
        db.execute(
            text(
                """
                UPDATE resources
                SET embedding_vec = CAST(:vec AS vector)
                WHERE id = :resource_id
                """
            ),
            {
                "vec": _embedding_to_pgvector_literal(embedding),
                "resource_id": resource_id,
            },
        )
    except Exception:  # noqa: BLE001
        # Keep JSON embedding as fallback when pgvector is unavailable.
        pass


def enrich_resource_with_ai(
    db: Session,
    resource: models.Resource,
    *,
    force: bool = False,
) -> bool:
    if not ai_service.is_enabled():
        return False
    if not force and not settings.AI_AUTO_ENRICH:
        return False
    if (
        not force
        and resource.ai_summary
        and resource.ai_tags
        and isinstance(resource.embedding_json, list)
        and resource.embedding_json
    ):
        return False

    source = extract_text_for_ai(resource)
    if len(source.strip()) < 10:
        return False

    try:
        summary, tags = ai_service.generate_summary_and_tags(
            source,
            title=resource.title,
            subject=resource.subject,
        )
        embedding = ai_service.generate_embedding(source)
    except ai_service.AIServiceError:
        return False

    resource.ai_summary = summary or None
    resource.ai_tags = tags or []
    resource.embedding_json = embedding or None
    resource.embedding_model = settings.AI_EMBEDDING_MODEL if embedding else None
    resource.ai_updated_at = datetime.now(timezone.utc)
    db.add(resource)
    _update_embedding_vector_column(db, resource_id=resource.id, embedding=embedding or None)
    db.commit()
    db.refresh(resource)
    return True


def enrich_resource_with_ai_background(resource_id: int, *, force: bool = False) -> None:
    from app.core.db_read_write import WriteSessionLocal

    db = WriteSessionLocal()
    try:
        resource = db.query(models.Resource).filter(models.Resource.id == resource_id).first()
        if not resource or resource.is_trashed:
            return
        enrich_resource_with_ai(db, resource, force=force)
    except Exception:  # noqa: BLE001
        logger.exception("Background AI enrich failed for resource_id=%s", resource_id)
    finally:
        db.close()


def delete_resource_binary(resource: models.Resource) -> None:
    if resource.storage_provider == models.StorageProvider.minio and resource.object_key:
        # Binary deletion now uses trash flow in trash_service.
        return

    if resource.file_path:
        filename = resource.file_path.split("/")[-1]
        local_path = UPLOAD_DIR / filename
        if local_path.exists():
            local_path.unlink()


def to_section_lite(section: models.ResourceSection | None) -> schemas.ResourceSectionLiteOut | None:
    if not section:
        return None
    return schemas.ResourceSectionLiteOut(id=section.id, code=section.code, name=section.name)


def to_resource_out(resource: models.Resource) -> schemas.ResourceOut:
    download_url = resolve_download_url(resource)
    section = resource.section

    return schemas.ResourceOut(
        id=resource.id,
        title=resource.title,
        description=resource.description,
        type=resource.type,
        subject=resource.subject,
        grade=resource.grade,
        tags=resource.tags,
        status=resource.status,
        resource_kind=resource.resource_kind,
        file_format=resource.file_format,
        difficulty=resource.difficulty,
        ai_summary=resource.ai_summary,
        ai_tags=resource.ai_tags or [],
        has_embedding=bool(resource.embedding_json),
        section_id=resource.section_id,
        section=to_section_lite(section),
        volume_code=resource.volume_code,
        source_filename=resource.source_filename,
        title_auto_generated=resource.title_auto_generated,
        rename_version=resource.rename_version,
        storage_provider=resource.storage_provider,
        object_key=resource.object_key,
        chapter_id=resource.chapter_id,
        chapter_ids=resolve_chapter_ids(resource),
        author_id=resource.author_id,
        reviewer_id=resource.reviewer_id,
        review_note=resource.review_note,
        file_path=resource.file_path,
        is_trashed=resource.is_trashed,
        trashed_at=resource.trashed_at,
        trashed_by=resource.trashed_by,
        trash_source=resource.trash_source,
        trash_has_binary=resource.trash_has_binary,
        download_url=download_url,
        preview_mode=resource.file_format,
        created_at=resource.created_at,
        updated_at=resource.updated_at,
    )


def _embedding_to_pgvector_literal(embedding: list[float]) -> str:
    return "[" + ",".join(f"{float(value):.8f}" for value in embedding) + "]"


def _semantic_candidate_resource_ids(
    db: Session,
    *,
    query: str,
    query_embedding: list[float] | None,
    candidate_limit: int,
) -> list[int]:
    cap = max(20, min(2000, int(candidate_limit)))
    seen: set[int] = set()
    ids: list[int] = []

    if settings.SEMANTIC_PGVECTOR_ENABLED and query_embedding:
        try:
            vector_literal = _embedding_to_pgvector_literal(query_embedding)
            rows = db.execute(
                text(
                    """
                    SELECT id
                    FROM resources
                    WHERE status = 'approved'
                      AND is_trashed = FALSE
                      AND embedding_vec IS NOT NULL
                    ORDER BY embedding_vec <=> CAST(:vec AS vector)
                    LIMIT :limit
                    """
                ),
                {"vec": vector_literal, "limit": cap},
            ).fetchall()
            for row in rows:
                rid = int(row[0])
                if rid not in seen:
                    seen.add(rid)
                    ids.append(rid)
        except Exception:  # noqa: BLE001
            # pgvector may be unavailable; lexical fallback still runs.
            pass

    lexical_cap = max(20, min(cap, 600))
    pattern = f"%{query.strip()}%"
    lexical_rows = (
        db.query(models.Resource.id)
        .filter(
            models.Resource.status == models.ResourceStatus.approved,
            models.Resource.is_trashed.is_(False),
            or_(
                models.Resource.title.ilike(pattern),
                models.Resource.description.ilike(pattern),
                models.Resource.ai_summary.ilike(pattern),
                func.array_to_string(models.Resource.tags, ",").ilike(pattern),
                func.array_to_string(models.Resource.ai_tags, ",").ilike(pattern),
                models.Resource.source_filename.ilike(pattern),
            ),
        )
        .order_by(models.Resource.updated_at.desc())
        .limit(lexical_cap)
        .all()
    )
    for row in lexical_rows:
        rid = int(row.id)
        if rid not in seen:
            seen.add(rid)
            ids.append(rid)

    if not ids:
        recent_rows = (
            db.query(models.Resource.id)
            .filter(
                models.Resource.status == models.ResourceStatus.approved,
                models.Resource.is_trashed.is_(False),
            )
            .order_by(models.Resource.updated_at.desc())
            .limit(cap)
            .all()
        )
        for row in recent_rows:
            rid = int(row.id)
            if rid not in seen:
                seen.add(rid)
                ids.append(rid)

    return ids[:cap]


def query_by_filters(
    db: Session,
    q: str | None,
    all_flag: bool,
    payload: dict | None,
    subject: str | None,
    grade: str | None,
    resource_kind: str | None,
    file_format: str | None,
    difficulty: str | None,
    chapter_id: int | None,
    section_id: int | None,
    status_filter: models.ResourceStatus | None,
):
    query = db.query(models.Resource).filter(models.Resource.is_trashed.is_(False))

    if q:
        pattern = f"%{q}%"
        query = query.filter(
            or_(
                models.Resource.title.ilike(pattern),
                models.Resource.description.ilike(pattern),
                models.Resource.ai_summary.ilike(pattern),
                func.array_to_string(models.Resource.tags, ",").ilike(pattern),
                func.array_to_string(models.Resource.ai_tags, ",").ilike(pattern),
                models.Resource.source_filename.ilike(pattern),
            )
        )

    if subject:
        query = query.filter(models.Resource.subject == subject)
    if grade:
        query = query.filter(models.Resource.grade == grade)
    if resource_kind:
        query = query.filter(models.Resource.resource_kind == resource_kind)
    if file_format:
        query = query.filter(models.Resource.file_format == file_format)
    if difficulty:
        query = query.filter(models.Resource.difficulty == difficulty)
    if section_id:
        query = query.filter(models.Resource.section_id == section_id)

    if chapter_id:
        linked_ids = select(models.ResourceChapterLink.resource_id).where(
            models.ResourceChapterLink.chapter_id == chapter_id
        )
        query = query.filter(
            or_(
                models.Resource.chapter_id == chapter_id,
                models.Resource.id.in_(linked_ids),
            )
        )

    if all_flag:
        role = payload.get("role") if payload else None
        if role != models.UserRole.admin.value:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Admin only",
            )
        if status_filter is not None:
            query = query.filter(models.Resource.status == status_filter)
    else:
        if status_filter is not None:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Admin only",
            )
        query = query.filter(models.Resource.status == models.ResourceStatus.approved)

    return query


def infer_section(
    db: Session,
    section_id: int | None,
    resource_kind: str | None,
    stage: str | None,
    subject: str | None,
) -> tuple[models.ResourceSection | None, str]:
    section = None
    if section_id is not None:
        section = db.query(models.ResourceSection).filter(models.ResourceSection.id == section_id).first()
        if not section:
            raise HTTPException(status_code=400, detail="Invalid section_id")
        if not section.is_enabled:
            raise HTTPException(status_code=400, detail="Section is disabled")

    stage_scope = stage or "senior"
    subject_scope = (subject or "").strip()

    if section is None and resource_kind:
        section = (
            db.query(models.ResourceSection)
            .filter(
                models.ResourceSection.stage == stage_scope,
                models.ResourceSection.subject == subject_scope,
                models.ResourceSection.code == resource_kind,
            )
            .first()
        )

    resolved_kind = resource_kind or "tutorial"
    if section:
        resolved_kind = section.code

    if resolved_kind not in schemas.RESOURCE_KIND_VALUES and not section:
        raise HTTPException(status_code=400, detail="Invalid resource_kind")

    return section, resolved_kind


@router.get("", response_model=schemas.ResourceListOut | list[schemas.ResourceOut])
def list_resources(
    q: str | None = Query(default=None, description="Search keyword"),
    all: bool = Query(default=False, description="Admin can list all resources"),
    subject: str | None = Query(default=None),
    grade: str | None = Query(default=None),
    resource_kind: str | None = Query(default=None),
    file_format: str | None = Query(default=None),
    difficulty: str | None = Query(default=None),
    chapter_id: int | None = Query(default=None),
    section_id: int | None = Query(default=None),
    status: models.ResourceStatus | None = Query(default=None),
    page: int = Query(default=1, ge=1),
    page_size: int = Query(default=20, ge=1, le=200),
    sort: str = Query(default="created_at_desc"),
    legacy_flat: bool = Query(default=False),
    db: Session = Depends(get_db_read),
    payload: dict | None = Depends(get_auth_payload_optional),
):
    query = query_by_filters(
        db,
        q,
        all,
        payload,
        subject,
        grade,
        resource_kind,
        file_format,
        difficulty,
        chapter_id,
        section_id,
        status,
    )
    order = models.Resource.created_at.desc()
    normalized_sort = (sort or "").strip().lower()
    if normalized_sort == "updated_at_desc":
        order = models.Resource.updated_at.desc()
    elif normalized_sort == "updated_at_asc":
        order = models.Resource.updated_at.asc()
    elif normalized_sort == "created_at_asc":
        order = models.Resource.created_at.asc()

    total = query.count()
    rows = (
        query.order_by(order)
        .offset((page - 1) * page_size)
        .limit(page_size)
        .all()
    )
    items = [to_resource_out(row) for row in rows]
    if legacy_flat:
        return items
    return schemas.ResourceListOut(
        items=items,
        total=total,
        page=page,
        page_size=page_size,
    )


@router.get("/ai-status", response_model=schemas.ResourceAiStatusOut)
def get_ai_status(
    _: models.User = Depends(get_current_user),
):
    return schemas.ResourceAiStatusOut(
        enabled=ai_service.is_enabled(),
        auto_enrich=bool(settings.AI_AUTO_ENRICH),
    )


@router.post("/semantic-search", response_model=schemas.SemanticSearchResponse)
def semantic_search(
    payload: schemas.SemanticSearchRequest,
    db: Session = Depends(get_db_read),
    _: models.User = Depends(get_current_user),
):
    query_embedding: list[float] | None = None
    if ai_service.is_enabled():
        try:
            query_embedding = ai_service.generate_embedding(payload.query)
        except ai_service.AIServiceError:
            query_embedding = None

    candidate_limit = payload.candidate_limit or settings.SEMANTIC_DEFAULT_CANDIDATE_LIMIT
    rerank_top_k = payload.rerank_top_k or settings.SEMANTIC_DEFAULT_RERANK_TOP_K

    candidate_ids = _semantic_candidate_resource_ids(
        db,
        query=payload.query,
        query_embedding=query_embedding,
        candidate_limit=candidate_limit,
    )
    if not candidate_ids:
        return schemas.SemanticSearchResponse(
            query=payload.query,
            answer=None,
            threshold=0.02,
            returned_count=0,
            scoring_profile="balanced_v1",
            results=[],
        )

    rows = (
        db.query(models.Resource)
        .filter(models.Resource.id.in_(candidate_ids))
        .order_by(models.Resource.updated_at.desc())
        .all()
    )
    row_map = {row.id: row for row in rows}
    ordered_rows = [row_map[item] for item in candidate_ids if item in row_map]

    candidates: list[semantic_ranker.SemanticCandidate] = []
    resource_map: dict[str, models.Resource] = {}
    for row in ordered_rows:
        candidate_id = f"resource:{row.id}"
        resource_map[candidate_id] = row
        candidates.append(
            semantic_ranker.SemanticCandidate(
                candidate_id=candidate_id,
                title=row.title or "",
                description=row.description or "",
                summary=row.ai_summary or "",
                tags=list(dict.fromkeys((row.ai_tags or []) + (row.tags or []))),
                embedding=row.embedding_json if isinstance(row.embedding_json, list) else None,
                chapter_id=row.chapter_id,
                section_id=row.section_id,
            )
        )

    ranked = semantic_ranker.rank_candidates(
        payload.query,
        candidates,
        query_embedding=query_embedding,
        top_k=max(1, min(20, payload.top_k, rerank_top_k)),
    )

    results: list[schemas.SemanticSearchItem] = []
    seen_canonical_keys: set[str] = set()
    for item in ranked.items:
        row = resource_map.get(item.candidate.candidate_id)
        if not row:
            continue
        canonical_key = resource_variants.build_canonical_key(
            resource_id=row.id,
            object_key=row.object_key,
        )
        if payload.dedupe and canonical_key in seen_canonical_keys:
            continue
        seen_canonical_keys.add(canonical_key)
        results.append(
            schemas.SemanticSearchItem(
                score=round(item.probability, 6),
                probability=round(item.probability, 6),
                factors=schemas.SemanticScoreFactorsOut(
                    vector=round(item.vector, 6),
                    summary=round(item.summary, 6),
                    content=round(item.content, 6),
                    tags=round(item.tags, 6),
                    raw=round(item.raw, 6),
                ),
                resource=to_resource_out(row),
                target=schemas.SemanticSearchTargetOut(
                    resource_id=row.id,
                    source_id=None,
                    canonical_key=canonical_key,
                    title=row.title,
                    file_format=row.file_format,
                    chapter_id=row.chapter_id,
                    section_id=row.section_id,
                    summary=row.ai_summary or row.description,
                    tags=list(dict.fromkeys((row.ai_tags or []) + (row.tags or []))),
                ),
                highlight_nodes=[],
                highlight_edges=[],
            )
        )

    answer = None
    if settings.AI_SEMANTIC_ENABLE_ANSWER and payload.include_answer and results and ai_service.is_enabled():
        contexts = []
        for item in results[:12]:
            row = resource_map.get(f"resource:{item.resource.id}") if item.resource else None
            if not row:
                continue
            contexts.append(
                {
                    "id": row.id,
                    "title": row.title,
                    "summary": row.ai_summary or row.description or "",
                    "snippet": (row.description or "")[:280],
                    "tags": row.ai_tags or row.tags or [],
                }
            )
        try:
            answer = ai_service.generate_rag_answer(payload.query, contexts)
        except ai_service.AIServiceError:
            answer = None

    return schemas.SemanticSearchResponse(
        query=payload.query,
        answer=answer,
        threshold=round(ranked.threshold, 6),
        returned_count=len(results),
        scoring_profile="balanced_v1",
        results=results,
    )


@router.get("/chapter/{chapter_id}/groups", response_model=schemas.ChapterGroupsOut)
def chapter_resources_groups(
    chapter_id: int,
    q: str | None = Query(default=None),
    file_format: str | None = Query(default=None),
    difficulty: str | None = Query(default=None),
    db: Session = Depends(get_db_read),
    _: models.User = Depends(get_current_user),
):
    chapter = db.query(models.Chapter).filter(models.Chapter.id == chapter_id).first()
    if not chapter:
        raise HTTPException(status_code=404, detail="Chapter not found")

    query = query_by_filters(
        db=db,
        q=q,
        all_flag=False,
        payload=None,
        subject=chapter.subject,
        grade=chapter.grade,
        resource_kind=None,
        file_format=file_format,
        difficulty=difficulty,
        chapter_id=chapter_id,
        section_id=None,
        status_filter=None,
    )
    rows = query.order_by(models.Resource.created_at.desc()).all()

    enabled_sections = (
        db.query(models.ResourceSection)
        .filter(
            models.ResourceSection.stage == chapter.stage,
            models.ResourceSection.subject == chapter.subject,
            models.ResourceSection.is_enabled.is_(True),
        )
        .order_by(models.ResourceSection.sort_order.asc(), models.ResourceSection.id.asc())
        .all()
    )
    enabled_ids = {row.id for row in enabled_sections}
    grouped: dict[int, list[schemas.ResourceOut]] = {row.id: [] for row in enabled_sections}
    unsectioned: list[schemas.ResourceOut] = []

    for row in rows:
        item = to_resource_out(row)
        if row.section_id and row.section_id in enabled_ids:
            grouped[row.section_id].append(item)
        else:
            unsectioned.append(item)

    groups: list[schemas.DynamicGroupOut] = []
    for section in enabled_sections:
        groups.append(
            schemas.DynamicGroupOut(
                section=schemas.ResourceSectionLiteOut(
                    id=section.id,
                    code=section.code,
                    name=section.name,
                ),
                items=grouped.get(section.id, []),
            )
        )

    if unsectioned:
        groups.append(schemas.DynamicGroupOut(section=None, items=unsectioned))

    return schemas.ChapterGroupsOut(chapter=schemas.ChapterOut.model_validate(chapter), groups=groups)


@router.get("/chapter/{chapter_id}", response_model=schemas.ChapterGroupsOut)
def chapter_resources_compat(
    chapter_id: int,
    q: str | None = Query(default=None),
    file_format: str | None = Query(default=None),
    difficulty: str | None = Query(default=None),
    db: Session = Depends(get_db_read),
    current_user: models.User = Depends(get_current_user),
):
    return chapter_resources_groups(
        chapter_id=chapter_id,
        q=q,
        file_format=file_format,
        difficulty=difficulty,
        db=db,
        _=current_user,
    )


@router.get("/mine", response_model=list[schemas.ResourceOut])
def my_resources(
    db: Session = Depends(get_db_write),
    current_user: models.User = Depends(get_current_user),
):
    rows = (
        db.query(models.Resource)
        .filter(
            models.Resource.author_id == current_user.id,
            models.Resource.is_trashed.is_(False),
        )
        .order_by(models.Resource.created_at.desc())
        .all()
    )
    return [to_resource_out(row) for row in rows]


@router.get("/upload-path-preview", response_model=schemas.UploadPathPreviewOut)
def upload_path_preview(
    filename: str = Query(..., min_length=1, max_length=255),
    chapter_id: int | None = Query(default=None),
    section_id: int | None = Query(default=None),
    volume_code: str | None = Query(default=None),
    low_confidence: bool = Query(default=False),
    db: Session = Depends(get_db_read),
    _: models.User = Depends(get_current_user),
):
    chapter_code = None
    section_code = None

    if chapter_id is not None:
        chapter = db.query(models.Chapter).filter(models.Chapter.id == chapter_id).first()
        if not chapter:
            raise HTTPException(status_code=400, detail="Invalid chapter_id")
        chapter_code = chapter.chapter_code
        volume_code = chapter.volume_code

    if section_id is not None:
        section = db.query(models.ResourceSection).filter(models.ResourceSection.id == section_id).first()
        if not section:
            raise HTTPException(status_code=400, detail="Invalid section_id")
        if not section.is_enabled:
            raise HTTPException(status_code=400, detail="Section is disabled")
        section_code = section.code

    object_key = build_resource_object_key(
        filename,
        chapter_code,
        section_code,
        volume_code=volume_code,
        low_confidence=low_confidence,
    )
    prefix = f"{Path(object_key).parent.as_posix()}/"
    return schemas.UploadPathPreviewOut(
        object_key=object_key,
        prefix=prefix,
        is_unassigned="/unassigned/" in object_key or object_key.startswith("resources/unassigned/"),
    )


def _to_auto_classify_response(result: chapter_classifier.ChapterClassification) -> schemas.AutoClassifyResponse:
    return schemas.AutoClassifyResponse(
        picked_chapter_id=result.chapter.id if result.chapter else None,
        picked_volume_code=result.volume_code,
        recommended_chapter_id=result.recommended_chapter_id,
        confidence=round(result.confidence, 6),
        confidence_level=result.confidence_level,  # type: ignore[arg-type]
        is_low_confidence=result.is_low_confidence,
        candidates=[
            schemas.AutoClassifyCandidateOut(
                chapter_id=item.chapter.id,
                volume_code=item.chapter.volume_code,
                title=f"{item.chapter.volume_name} {item.chapter.chapter_code} {item.chapter.title}",
                score=round(item.probability, 6),
                probability=round(item.probability, 6),
                reasons=item.reasons,
                rule_score=round(item.rule_score, 6),
                lexical_score=round(item.lexical_score, 6),
                vector_score=round(item.vector_score, 6),
                final_score=round(item.final_score, 6),
            )
            for item in result.candidates
        ],
        rule_hits=result.rule_hits,
        catalog_version="pep2019_v1",
        reason=result.reason,
    )


def _validate_external_url(raw_url: str) -> str:
    value = (raw_url or "").strip()
    if not value:
        raise HTTPException(status_code=400, detail="external_url is required")
    try:
        parsed = urlparse(value)
    except ValueError as error:
        raise HTTPException(status_code=400, detail="Invalid external_url") from error
    if parsed.scheme not in {"http", "https"} or not parsed.netloc:
        raise HTTPException(status_code=400, detail="external_url must be http/https")
    return value


@router.post("/auto-classify", response_model=schemas.AutoClassifyResponse)
def auto_classify_resource(
    title: str = Form(default=""),
    description: str = Form(default=""),
    tags: str = Form(default=""),
    subject: str = Form(default="物理"),
    stage: str = Form(default="senior"),
    volume_code: str = Form(default=""),
    selected_volume_code: str = Form(default=""),
    external_url: str = Form(default=""),
    file: UploadFile | None = File(default=None),
    db: Session = Depends(get_db_write),
    _: models.User = Depends(get_current_user),
):
    parsed_tags = parse_tags(tags)
    normalized_url = external_url.strip()
    if normalized_url:
        normalized_url = _validate_external_url(normalized_url)
    file_format = detect_file_format(file.filename if file else None)
    preview_text = _read_upload_preview_text(file, file_format)
    result = chapter_classifier.classify_chapter(
        db,
        stage=stage.strip() or "senior",
        subject=subject.strip() or "物理",
        title=title.strip(),
        description=description.strip(),
        tags=parsed_tags,
        filename=file.filename if file else "",
        external_url=normalized_url,
        content_text=preview_text,
        volume_code=selected_volume_code.strip() or volume_code.strip() or None,
        top_k=3,
    )
    return _to_auto_classify_response(result)


@router.get("/pending", response_model=list[schemas.ResourceOut])
def pending_resources(
    db: Session = Depends(get_db_write),
    _: models.User = Depends(get_current_admin),
):
    rows = (
        db.query(models.Resource)
        .filter(
            models.Resource.status == models.ResourceStatus.pending,
            models.Resource.is_trashed.is_(False),
        )
        .order_by(models.Resource.created_at.asc())
        .all()
    )
    return [to_resource_out(row) for row in rows]


@router.post("/text-to-md", response_model=schemas.TextToMarkdownResponse)
def text_to_markdown(
    payload: schemas.TextToMarkdownRequest,
    _: models.User = Depends(get_current_user),
):
    try:
        markdown = text_to_markdown_with_mineru(payload.text, payload.title)
    except MinerUAPIError as error:
        raise HTTPException(status_code=502, detail=str(error)) from error

    return schemas.TextToMarkdownResponse(markdown=markdown)


@router.post(
    "",
    response_model=schemas.ResourceOut,
    status_code=status.HTTP_201_CREATED,
)
def create_resource(
    background_tasks: BackgroundTasks,
    title: str = Form(default=""),
    type: str = Form(...),
    description: str = Form(default=""),
    subject: str = Form(default="物理"),
    grade: str = Form(default=""),
    tags: str = Form(default=""),
    resource_kind: str = Form(default="tutorial"),
    difficulty: str = Form(default=""),
    chapter_id: int | None = Form(default=None),
    section_id: int | None = Form(default=None),
    volume_code: str = Form(default=""),
    file: UploadFile | None = File(default=None),
    db: Session = Depends(get_db_write),
    current_user: models.User = Depends(get_current_user),
):
    parsed_tags = parse_tags(tags)
    file_format = detect_file_format(file.filename if file else None)
    chapter = None
    resolved_volume_code = volume_code.strip() or None
    if chapter_id is not None:
        chapter = db.query(models.Chapter).filter(models.Chapter.id == chapter_id).first()
        if not chapter:
            raise HTTPException(status_code=400, detail="Invalid chapter_id")
        if not chapter.is_enabled:
            raise HTTPException(status_code=400, detail="Chapter is disabled")
        resolved_volume_code = chapter.volume_code

    resolved_subject = (subject.strip() or (chapter.subject if chapter else "物理")).strip()
    if not resolved_subject:
        resolved_subject = "物理"
    resolved_grade = (grade.strip() or (chapter.grade if chapter else "")).strip()
    stage = chapter.stage if chapter else "senior"
    section, resolved_kind = infer_section(
        db=db,
        section_id=section_id,
        resource_kind=resource_kind,
        stage=stage,
        subject=resolved_subject,
    )

    auto_result: chapter_classifier.ChapterClassification | None = None
    if chapter is None:
        auto_result = chapter_classifier.classify_chapter(
            db,
            stage=stage,
            subject=resolved_subject,
            title=title.strip(),
            description=description.strip(),
            tags=parsed_tags,
            filename=file.filename if file else "",
            content_text=_read_upload_preview_text(file, file_format),
            volume_code=resolved_volume_code,
            top_k=3,
        )
        if auto_result.chapter and not auto_result.is_low_confidence:
            chapter = auto_result.chapter
            chapter_id = chapter.id
            resolved_volume_code = chapter.volume_code
            if not resolved_grade:
                resolved_grade = chapter.grade
        else:
            resolved_volume_code = resolved_volume_code or auto_result.volume_code

    if stage == "senior" and resolved_subject == "物理" and chapter is None:
        raise HTTPException(status_code=400, detail="高中物理资源必须先选择人教版目录章节")

    resolved_title, clean_base_name, title_auto_generated = compose_resource_name(
        raw_title=title,
        chapter=chapter,
        section=section,
        filename=file.filename if file else None,
        tags=parsed_tags,
        description=description,
        volume_code_override=resolved_volume_code,
    )

    object_key = None
    storage_provider = models.StorageProvider.local
    file_path = None
    source_filename = file.filename if file and file.filename else None

    if file and file.filename:
        try:
            planned_key = build_resource_object_key(
                file.filename,
                chapter.chapter_code if chapter else None,
                section.code if section else None,
                volume_code=resolved_volume_code,
                base_name=clean_base_name,
                low_confidence=chapter is None and bool(resolved_volume_code),
            )
            object_key, _ = upload_file(file, object_key=planned_key)
            storage_provider = models.StorageProvider.minio
        except Exception:  # noqa: BLE001
            file_suffix = Path(file.filename).suffix.lower()
            filename = f"{uuid4().hex}{file_suffix}"
            save_path = UPLOAD_DIR / filename
            with save_path.open("wb") as destination:
                shutil.copyfileobj(file.file, destination)
            file_path = f"/uploads/{filename}"
            storage_provider = models.StorageProvider.local

    resource = models.Resource(
        title=resolved_title,
        description=description.strip() or None,
        type=type.strip(),
        subject=resolved_subject or None,
        grade=resolved_grade or None,
        tags=parsed_tags,
        status=models.ResourceStatus.pending,
        resource_kind=resolved_kind,
        file_format=file_format,
        difficulty=difficulty.strip() or None,
        section_id=section.id if section else None,
        volume_code=resolved_volume_code,
        source_filename=source_filename,
        title_auto_generated=title_auto_generated,
        rename_version="v1",
        storage_provider=storage_provider,
        object_key=object_key,
        chapter_id=chapter.id if chapter else None,
        author_id=current_user.id,
        file_path=file_path,
    )
    db.add(resource)
    db.commit()
    db.refresh(resource)

    if resource.storage_provider == models.StorageProvider.minio and resource.object_key:
        resource_variants.ensure_resource_origin_variant(db, resource)
        db.commit()

    if chapter:
        link = models.ResourceChapterLink(resource_id=resource.id, chapter_id=chapter.id)
        db.add(link)
        db.commit()
        db.refresh(resource)

    if resource.storage_provider == models.StorageProvider.minio and resource.object_key:
        try:
            if resource.file_format == "ppt":
                preview_key = ensure_presentation_pdf_preview(resource.object_key, force=False)
                if preview_key:
                    resource_variants.ensure_resource_preview_pdf_variant(
                        db,
                        resource=resource,
                        preview_key=preview_key,
                    )
                    db.commit()
            else:
                preview_key = ensure_legacy_pdf_preview(resource.object_key, force=False)
                if preview_key:
                    resource_variants.ensure_resource_preview_pdf_variant(
                        db,
                        resource=resource,
                        preview_key=preview_key,
                    )
                    db.commit()
        except Exception:  # noqa: BLE001
            pass

    background_tasks.add_task(enrich_resource_with_ai_background, resource.id, force=False)

    return to_resource_out(resource)


@router.post("/{resource_id}/index/chapter", response_model=schemas.ResourceOut)
def link_resource_chapter(
    resource_id: int,
    payload: schemas.LinkChapterRequest,
    db: Session = Depends(get_db_write),
    current_user: models.User = Depends(get_current_user),
):
    resource = db.query(models.Resource).filter(models.Resource.id == resource_id).first()
    if not resource:
        raise HTTPException(status_code=404, detail="Resource not found")
    if resource.is_trashed:
        raise HTTPException(status_code=404, detail="Resource not available")
    if resource.is_trashed:
        raise HTTPException(status_code=404, detail="Resource not available")

    if current_user.role != models.UserRole.admin and resource.author_id != current_user.id:
        raise HTTPException(status_code=403, detail="No permission")

    chapter = db.query(models.Chapter).filter(models.Chapter.id == payload.chapter_id).first()
    if not chapter:
        raise HTTPException(status_code=400, detail="Invalid chapter_id")

    exists = (
        db.query(models.ResourceChapterLink)
        .filter(
            models.ResourceChapterLink.resource_id == resource_id,
            models.ResourceChapterLink.chapter_id == payload.chapter_id,
        )
        .first()
    )
    if not exists:
        db.add(
            models.ResourceChapterLink(
                resource_id=resource_id,
                chapter_id=payload.chapter_id,
            )
        )
        db.commit()

    db.refresh(resource)
    return to_resource_out(resource)


@router.get("/{resource_id}/office-config", response_model=schemas.OfficeConfigOut)
def resource_office_config(
    resource_id: int,
    db: Session = Depends(get_db_read),
    current_user: models.User = Depends(get_current_user),
):
    if not settings.ONLYOFFICE_ENABLED:
        raise HTTPException(status_code=503, detail="OnlyOffice is disabled")

    resource = db.query(models.Resource).filter(models.Resource.id == resource_id).first()
    if not resource:
        raise HTTPException(status_code=404, detail="Resource not found")

    if (
        resource.status != models.ResourceStatus.approved
        and current_user.role != models.UserRole.admin
        and resource.author_id != current_user.id
    ):
        raise HTTPException(status_code=403, detail="Resource not approved")

    if not is_office_resource(resource):
        raise HTTPException(status_code=400, detail="Resource is not an office file")
    if is_legacy_office_resource(resource):
        raise HTTPException(status_code=400, detail="Legacy office format uses PDF preview")
    if resource.storage_provider != models.StorageProvider.minio or not resource.object_key:
        raise HTTPException(status_code=400, detail="Resource object is not available in MinIO")

    editable = current_user.role == models.UserRole.admin
    return build_office_config(
        object_key=resource.object_key,
        filename=resolve_resource_filename(resource) or resource.title,
        current_user=current_user,
        editable=editable,
    )


@router.get("/{resource_id}/access-urls", response_model=schemas.AccessUrlsOut)
def resource_access_urls(
    resource_id: int,
    db: Session = Depends(get_db_read),
    current_user: models.User = Depends(get_current_user),
):
    resource = db.query(models.Resource).filter(models.Resource.id == resource_id).first()
    if not resource:
        raise HTTPException(status_code=404, detail="Resource not found")
    if resource.is_trashed:
        raise HTTPException(status_code=404, detail="Resource not available")

    if (
        resource.status != models.ResourceStatus.approved
        and current_user.role != models.UserRole.admin
        and resource.author_id != current_user.id
    ):
        raise HTTPException(status_code=403, detail="Resource not approved")

    access = resolve_access_urls(resource, user_id=current_user.id)
    if not access:
        raise HTTPException(status_code=404, detail="Resource attachment not found")
    return access


@router.get("/{resource_id}/preview", response_model=schemas.ResourcePreviewOut)
def resource_preview(
    resource_id: int,
    db: Session = Depends(get_db_read),
    current_user: models.User = Depends(get_current_user),
):
    resource = db.query(models.Resource).filter(models.Resource.id == resource_id).first()
    if not resource:
        raise HTTPException(status_code=404, detail="Resource not found")
    if resource.is_trashed:
        raise HTTPException(status_code=404, detail="Resource not available")

    if resource.status != models.ResourceStatus.approved:
        raise HTTPException(status_code=403, detail="Resource not approved")

    access = resolve_access_urls(resource, user_id=current_user.id)
    open_url = access.open_url if access else None
    download_url = access.download_url if access else None
    mode = resource.file_format or "other"
    suffix = resolve_resource_suffix(resource)

    if mode == "ppt":
        presentation_pdf_key = resolve_presentation_preview_key(resource)
        if presentation_pdf_key:
            preview_open_url, preview_download_url = build_storage_access_urls(
                object_key=presentation_pdf_key,
                user_id=current_user.id,
            )
            return schemas.ResourcePreviewOut(
                mode="pdf",
                url=preview_open_url,
                open_url=preview_open_url,
                download_url=preview_download_url,
            )
        if is_legacy_office_suffix(suffix):
            legacy_pdf_key = resolve_legacy_preview_key(resource)
            if legacy_pdf_key:
                preview_open_url, preview_download_url = build_storage_access_urls(
                    object_key=legacy_pdf_key,
                    user_id=current_user.id,
                )
                return schemas.ResourcePreviewOut(
                    mode="pdf",
                    url=preview_open_url,
                    open_url=preview_open_url,
                    download_url=preview_download_url,
                )
        return schemas.ResourcePreviewOut(
            mode="other",
            url=open_url,
            open_url=open_url,
            download_url=download_url,
            content="PPT 预览转换失败，请稍后重试或下载查看。",
        )

    if is_office_resource(resource) and resource.storage_provider == models.StorageProvider.minio and resource.object_key:
        return schemas.ResourcePreviewOut(
            mode=mode,
            url=open_url,
            open_url=open_url,
            download_url=download_url,
        )

    if mode in {"pdf", "video", "image", "audio"}:
        return schemas.ResourcePreviewOut(
            mode=mode,
            url=open_url,
            open_url=open_url,
            download_url=download_url,
        )

    if mode == "html":
        repaired_html: str | None = None
        try:
            raw = read_preview_bytes(resource)
            repaired_html = repair_html_preview(raw.decode("utf-8", errors="ignore"))
        except HTTPException:
            repaired_html = None
        except Exception:  # noqa: BLE001
            repaired_html = None
        return schemas.ResourcePreviewOut(
            mode=mode,
            url=open_url,
            open_url=open_url,
            download_url=download_url,
            content=repaired_html,
        )

    if mode in {"markdown", "word", "excel", "ppt"}:
        raw = read_preview_bytes(resource)
        try:
            if mode == "markdown":
                content = raw.decode("utf-8", errors="ignore")
            elif mode == "word":
                content = render_word_preview(raw, suffix)
            elif mode == "excel":
                content = render_excel_preview(raw, suffix)
            else:
                content = render_ppt_preview(raw, suffix)
        except Exception:  # noqa: BLE001
            content = "预览解析失败，请下载后查看完整内容。"
        return schemas.ResourcePreviewOut(
            mode=mode,
            url=open_url,
            open_url=open_url,
            download_url=download_url,
            content=content,
        )

    return schemas.ResourcePreviewOut(
        mode="other",
        url=open_url,
        open_url=open_url,
        download_url=download_url,
    )


@router.patch("/{resource_id}/review", response_model=schemas.ResourceOut)
def review_resource(
    resource_id: int,
    payload: schemas.ReviewRequest,
    db: Session = Depends(get_db_write),
    current_admin: models.User = Depends(get_current_admin),
):
    resource = db.query(models.Resource).filter(models.Resource.id == resource_id).first()
    if not resource:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Resource not found",
        )
    if resource.is_trashed:
        raise HTTPException(status_code=400, detail="Resource is trashed")

    if payload.status in {models.ResourceStatus.pending, models.ResourceStatus.hidden}:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Invalid review status",
        )

    resource.status = payload.status
    resource.review_note = payload.review_note
    resource.reviewer_id = current_admin.id

    rag_sync.sync_resource_to_workspaces(
        db,
        [resource.id],
        actor_id=current_admin.id,
        reason="resource_review",
    )
    db.commit()
    db.refresh(resource)

    if payload.status == models.ResourceStatus.approved:
        enrich_resource_with_ai(db, resource)
        rag_sync.sync_resource_to_workspaces(
            db,
            [resource.id],
            actor_id=current_admin.id,
            reason="resource_review_ai_refresh",
        )
        db.commit()
        db.refresh(resource)

    return to_resource_out(resource)


@router.patch("/{resource_id}/visibility", response_model=schemas.ResourceOut)
def set_resource_visibility(
    resource_id: int,
    payload: schemas.ResourceVisibilityRequest,
    db: Session = Depends(get_db_write),
    current_admin: models.User = Depends(get_current_admin),
):
    resource = db.query(models.Resource).filter(models.Resource.id == resource_id).first()
    if not resource:
        raise HTTPException(status_code=404, detail="Resource not found")
    if resource.is_trashed:
        raise HTTPException(status_code=400, detail="Resource is trashed")

    target_status = (
        models.ResourceStatus.approved
        if payload.visibility == "public"
        else models.ResourceStatus.hidden
    )
    resource.status = target_status
    resource.reviewer_id = current_admin.id
    if payload.note and payload.note.strip():
        resource.review_note = payload.note.strip()
    db.add(resource)

    rag_sync.sync_resource_to_workspaces(
        db,
        [resource.id],
        actor_id=current_admin.id,
        reason=f"resource_visibility_{payload.visibility}",
    )
    db.commit()
    db.refresh(resource)
    return to_resource_out(resource)


@router.post("/bulk-manage", response_model=schemas.ResourceBulkManageOut)
def bulk_manage_resources(
    payload: schemas.ResourceBulkManageRequest,
    db: Session = Depends(get_db_write),
    current_admin: models.User = Depends(get_current_admin),
):
    requested_ids = list(dict.fromkeys([item for item in payload.resource_ids if item > 0]))
    if not requested_ids:
        return schemas.ResourceBulkManageOut(
            action=payload.action,
            requested=0,
            succeeded=0,
            failed=0,
            errors=[],
        )

    rows = (
        db.query(models.Resource)
        .filter(models.Resource.id.in_(requested_ids))
        .all()
    )
    resource_map = {row.id: row for row in rows}

    succeeded_ids: list[int] = []
    errors: list[schemas.ResourceBulkManageErrorOut] = []

    for resource_id in requested_ids:
        resource = resource_map.get(resource_id)
        if not resource:
            errors.append(schemas.ResourceBulkManageErrorOut(resource_id=resource_id, reason="Resource not found"))
            continue

        if payload.action in {"publish", "hide"}:
            if resource.is_trashed:
                errors.append(schemas.ResourceBulkManageErrorOut(resource_id=resource_id, reason="Resource is trashed"))
                continue
            target_status = (
                models.ResourceStatus.approved
                if payload.action == "publish"
                else models.ResourceStatus.hidden
            )
            resource.status = target_status
            resource.reviewer_id = current_admin.id
            if payload.note and payload.note.strip():
                resource.review_note = payload.note.strip()
            db.add(resource)
            succeeded_ids.append(resource_id)
            continue

        if payload.action == "trash":
            if resource.is_trashed:
                succeeded_ids.append(resource_id)
                continue
            trash_service.trash_resource(
                db,
                resource,
                source=trash_service.TRASH_SOURCE_RESOURCE_API,
                deleted_by=current_admin.id,
                scope=trash_service.TRASH_SCOPE_RESOURCE,
                meta={"from": "resources.bulk_manage"},
            )
            succeeded_ids.append(resource_id)
            continue

        errors.append(schemas.ResourceBulkManageErrorOut(resource_id=resource_id, reason="Unsupported action"))

    if succeeded_ids:
        rag_sync.sync_resource_to_workspaces(
            db,
            succeeded_ids,
            actor_id=current_admin.id,
            reason=f"resource_bulk_{payload.action}",
        )

    db.commit()

    return schemas.ResourceBulkManageOut(
        action=payload.action,
        requested=len(requested_ids),
        succeeded=len(succeeded_ids),
        failed=len(errors),
        errors=errors,
    )


@router.patch("/{resource_id}/tags", response_model=schemas.ResourceOut)
def update_resource_tags(
    resource_id: int,
    payload: schemas.ResourceTagsUpdateRequest,
    db: Session = Depends(get_db_write),
    current_admin: models.User = Depends(get_current_admin),
):
    resource = db.query(models.Resource).filter(models.Resource.id == resource_id).first()
    if not resource:
        raise HTTPException(status_code=404, detail="Resource not found")
    if resource.is_trashed:
        raise HTTPException(status_code=400, detail="Resource is trashed")

    incoming_tags = parse_tag_list(payload.tags)
    current_tags = parse_tag_list(resource.tags or [])
    if payload.mode == "append":
        resource.tags = list(dict.fromkeys(current_tags + incoming_tags))
    else:
        resource.tags = incoming_tags

    db.add(resource)
    rag_sync.sync_resource_to_workspaces(
        db,
        [resource.id],
        actor_id=current_admin.id,
        reason=f"resource_tags_{payload.mode}",
    )
    db.commit()
    db.refresh(resource)
    return to_resource_out(resource)


@router.post("/{resource_id}/tags/adopt-ai", response_model=schemas.ResourceOut)
def adopt_ai_tags(
    resource_id: int,
    payload: schemas.ResourceTagsAdoptAIRequest,
    db: Session = Depends(get_db_write),
    current_admin: models.User = Depends(get_current_admin),
):
    resource = db.query(models.Resource).filter(models.Resource.id == resource_id).first()
    if not resource:
        raise HTTPException(status_code=404, detail="Resource not found")
    if resource.is_trashed:
        raise HTTPException(status_code=400, detail="Resource is trashed")

    ai_tags = parse_tag_list(resource.ai_tags or [])
    current_tags = parse_tag_list(resource.tags or [])
    if payload.strategy == "replace":
        resource.tags = ai_tags
    else:
        resource.tags = list(dict.fromkeys(current_tags + ai_tags))

    db.add(resource)
    rag_sync.sync_resource_to_workspaces(
        db,
        [resource.id],
        actor_id=current_admin.id,
        reason=f"resource_tags_adopt_{payload.strategy}",
    )
    db.commit()
    db.refresh(resource)
    return to_resource_out(resource)


@router.post("/{resource_id}/ai-enrich", response_model=schemas.ResourceOut)
def manual_ai_enrich(
    resource_id: int,
    db: Session = Depends(get_db_write),
    _: models.User = Depends(get_current_admin),
):
    if not ai_service.is_enabled():
        raise HTTPException(status_code=503, detail="AI service is not configured")

    resource = db.query(models.Resource).filter(models.Resource.id == resource_id).first()
    if not resource:
        raise HTTPException(status_code=404, detail="Resource not found")
    if resource.is_trashed:
        raise HTTPException(status_code=400, detail="Resource is trashed")

    changed = enrich_resource_with_ai(db, resource, force=True)
    if not changed:
        raise HTTPException(status_code=400, detail="Unable to enrich resource by AI")
    return to_resource_out(resource)


@router.post("/ai-reindex")
def ai_reindex(
    limit: int = Query(default=50, ge=1, le=500),
    db: Session = Depends(get_db_write),
    _: models.User = Depends(get_current_admin),
):
    if not ai_service.is_enabled():
        raise HTTPException(status_code=503, detail="AI service is not configured")

    rows = (
        db.query(models.Resource)
        .filter(
            models.Resource.status == models.ResourceStatus.approved,
            models.Resource.is_trashed.is_(False),
        )
        .order_by(models.Resource.updated_at.desc())
        .limit(limit)
        .all()
    )

    done = 0
    skipped = 0
    for row in rows:
        changed = enrich_resource_with_ai(db, row, force=False)
        if changed:
            done += 1
        else:
            skipped += 1

    return {"processed": len(rows), "enriched": done, "skipped": skipped}


@router.delete("/{resource_id}", status_code=status.HTTP_204_NO_CONTENT)
def delete_resource(
    resource_id: int,
    db: Session = Depends(get_db_write),
    current_admin: models.User = Depends(get_current_admin),
):
    resource = db.query(models.Resource).filter(models.Resource.id == resource_id).first()
    if not resource:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Resource not found",
        )

    trash_service.trash_resource(
        db,
        resource,
        source=trash_service.TRASH_SOURCE_RESOURCE_API,
        deleted_by=current_admin.id,
        scope=trash_service.TRASH_SCOPE_RESOURCE,
        meta={"from": "resources.delete"},
    )
    rag_sync.sync_resource_to_workspaces(
        db,
        [resource.id],
        actor_id=current_admin.id,
        reason="resource_delete",
    )
    db.commit()
